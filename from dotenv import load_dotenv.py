from dotenv import load_dotenv
load_dotenv()

import os
import json
import tempfile
import subprocess
from datetime import datetime
from flask import Flask, request, jsonify, send_file, render_template
import anthropic
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

app = Flask(__name__)
client = anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY"))

FFMPEG = "/usr/local/bin/ffmpeg"
RECORDINGS_DIR = os.path.join(os.path.dirname(__file__), "recordings")
TRANSCRIPTS_DIR = os.path.join(os.path.dirname(__file__), "transcripts")
os.makedirs(RECORDINGS_DIR, exist_ok=True)
os.makedirs(TRANSCRIPTS_DIR, exist_ok=True)

PALETTE = {
    "primary":   "1E3A5F",
    "accent":    "2E86AB",
    "light":     "A8DADC",
    "white":     "FFFFFF",
    "bg_light":  "F4F8FB",
    "text_dark": "1A1A2E",
    "text_mid":  "4A5568",
}

def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def transcribe_audio(audio_bytes: bytes) -> str:
    from groq import Groq
    from datetime import datetime

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

    # Save raw audio
    audio_path = os.path.join(RECORDINGS_DIR, f"recording_{timestamp}.webm")
    with open(audio_path, "wb") as f:
        f.write(audio_bytes)
    print(f"✅ Audio saved: {audio_path}")

    try:
        groq_client = Groq(api_key=os.environ.get("GROQ_API_KEY"))

        with open(audio_path, "rb") as f:
            result = groq_client.audio.transcriptions.create(
                file=(f"recording_{timestamp}.webm", f),
                model="whisper-large-v3",
                response_format="text"
            )

        transcript = result.strip() if isinstance(result, str) else result.text.strip()

        # Save transcript
        txt_path = os.path.join(TRANSCRIPTS_DIR, f"transcript_{timestamp}.txt")
        with open(txt_path, "w") as f:
            f.write(f"Timestamp: {timestamp}\n")
            f.write(f"Audio file: recording_{timestamp}.webm\n")
            f.write("-" * 40 + "\n")
            f.write(transcript)
        print(f"✅ Transcript saved: {txt_path}")

        return transcript
        groq_client = Groq(api_key=os.environ.get("GROQ_API_KEY"))

        with open(audio_path, "rb") as f:
            result = groq_client.audio.transcriptions.create(
                file=(f"recording_{timestamp}.webm", f),
                model="whisper-large-v3",
                response_format="text"
            )

        transcript = result.strip() if isinstance(result, str) else result.text.strip()

        # Save transcript
        txt_path = os.path.join(TRANSCRIPTS_DIR, f"transcript_{timestamp}.txt")
        with open(txt_path, "w") as f:
            f.write(f"Timestamp: {timestamp}\n")
            f.write(f"Audio file: recording_{timestamp}.webm\n")
            f.write("-" * 40 + "\n")
            f.write(transcript)
        print(f"✅ Transcript saved: {txt_path}")

        return transcript

    except Exception as e:
        raise RuntimeError(f"Transcription failed: {e}")

def build_slides_from_json(slides_data):
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for i, info in enumerate(slides_data):
        slide = prs.slides.add_slide(blank)
        is_title = (i == 0)

        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = hex_to_rgb(PALETTE["primary"] if is_title else PALETTE["bg_light"])

        if not is_title:
            r = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.08))
            r.fill.solid()
            r.fill.fore_color.rgb = hex_to_rgb(PALETTE["accent"])
            r.line.fill.background()

        title = info.get("title", "")
        if is_title:
            tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.33), Inches(1.8))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = title
            run.font.size = Pt(44)
            run.font.bold = True
            run.font.color.rgb = hex_to_rgb(PALETTE["white"])
            sub = info.get("subtitle", "")
            if sub:
                sb = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.33), Inches(0.8))
                sp = sb.text_frame.paragraphs[0]
                sp.alignment = PP_ALIGN.CENTER
                sr2 = sp.add_run()
                sr2.text = sub
                sr2.font.size = Pt(20)
                sr2.font.color.rgb = hex_to_rgb(PALETTE["light"])
        else:
            hdr = slide.shapes.add_shape(1, Inches(0), Inches(0.08), Inches(13.33), Inches(1.1))
            hdr.fill.solid()
            hdr.fill.fore_color.rgb = hex_to_rgb(PALETTE["primary"])
            hdr.line.fill.background()
            tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.95))
            tf = tb.text_frame
            tf.word_wrap = True
            tp = tf.paragraphs[0]
            tp.alignment = PP_ALIGN.LEFT
            tr = tp.add_run()
            tr.text = title
            tr.font.size = Pt(28)
            tr.font.bold = True
            tr.font.color.rgb = hex_to_rgb(PALETTE["white"])

        bullets = info.get("bullets", [])
        if bullets and not is_title:
            cb = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.13), Inches(5.5))
            cf = cb.text_frame
            cf.word_wrap = True
            for idx, b in enumerate(bullets):
                para = cf.paragraphs[0] if idx == 0 else cf.add_paragraph()
                para.space_before = Pt(6)
                run = para.add_run()
                run.text = f"▸  {b}"
                run.font.size = Pt(20)
                run.font.color.rgb = hex_to_rgb(PALETTE["text_dark"])

        if not is_title:
            nb = slide.shapes.add_textbox(Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3))
            np = nb.text_frame.paragraphs[0]
            np.alignment = PP_ALIGN.RIGHT
            nr = np.add_run()
            nr.text = str(i)
            nr.font.size = Pt(11)
            nr.font.color.rgb = hex_to_rgb(PALETTE["text_mid"])

    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(tmp.name)
    return tmp.name

def ask_claude_for_slides(transcript):
    system = """You are a presentation expert. Given a voice transcript, extract key ideas and structure them as slides.
Return ONLY a valid JSON array. No markdown, no explanation. Format:
[
  {"title": "Presentation Title", "subtitle": "Optional subtitle", "type": "title"},
  {"title": "Slide Title", "bullets": ["Point one", "Point two", "Point three"], "type": "content"}
]
Rules:
- First slide must be title slide with "type": "title"
- Each content slide has 3-5 bullet points
- Aim for 4-7 slides total
- Keep titles under 8 words
- Bullets should be clear sentences"""

    msg = client.messages.create(
        model="claude-sonnet-4-6", max_tokens=2000, system=system,
        messages=[{"role": "user", "content": f"Convert this transcript into slides:\n\n{transcript}"}]
    )
    raw = msg.content[0].text.strip()
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
    return json.loads(raw.strip())

@app.route("/")
def index():
    return render_template("index.html")

@app.route("/transcribe", methods=["POST"])
def transcribe():
    if "audio" not in request.files:
        return jsonify({"error": "No audio file provided"}), 400
    try:
        text = transcribe_audio(request.files["audio"].read())
        return jsonify({"transcript": text})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/generate", methods=["POST"])
def generate():
    data = request.get_json()
    transcript = data.get("transcript", "").strip()
    if not transcript:
        return jsonify({"error": "Transcript is empty"}), 400
    try:
        return jsonify({"slides": ask_claude_for_slides(transcript)})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/download", methods=["POST"])
def download():
    slides = request.get_json().get("slides", [])
    if not slides:
        return jsonify({"error": "No slide data"}), 400
    try:
        path = build_slides_from_json(slides)
        return send_file(path, as_attachment=True, download_name="presentation.pptx",
                         mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/files", methods=["GET"])
def list_files():
    recordings = sorted(os.listdir(RECORDINGS_DIR), reverse=True)
    transcripts = sorted(os.listdir(TRANSCRIPTS_DIR), reverse=True)
    return jsonify({"recordings": recordings, "transcripts": transcripts})

@app.route("/transcripts/<filename>", methods=["GET"])
def get_transcript(filename):
    path = os.path.join(TRANSCRIPTS_DIR, filename)
    if not os.path.exists(path):
        return jsonify({"error": "File not found"}), 404
    with open(path) as f:
        return jsonify({"content": f.read()})

if __name__ == "__main__":
    print("🎙️  Voice-to-Slides running at http://localhost:5000")
    app.run(debug=True, port=5000)